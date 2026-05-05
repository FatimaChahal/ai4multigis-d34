"""
T3.4 Conclusion - Slide 15
Generates a PowerPoint slide with 3 takeaway blocks + central statement

Requirements:
    pip install python-pptx

Run:
    python slide15_conclusion.py
Output:
    slide15_conclusion.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Color palette ──────────────────────────────────────────────────────────────
DARK_BG     = RGBColor(0x06, 0x2A, 0x40)   # slide background (deep navy)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
MID_GRAY    = RGBColor(0x8A, 0x99, 0xA8)   # muted text
CARD_BG     = RGBColor(0x0C, 0x3A, 0x52)   # card background

BLUE        = RGBColor(0x17, 0x8A, 0xDD)   # block 1 — what we built
GREEN       = RGBColor(0x1D, 0x9E, 0x75)   # block 2 — what we proved
TEAL        = RGBColor(0x02, 0x80, 0x90)   # block 3 — what it means

ACCENT_LINE = RGBColor(0x17, 0x8A, 0xDD)   # separator line

# ── Slide dimensions (16:9) ───────────────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_text_lines(slide, lines, left, top, width, height,
                   font_size=12, color=WHITE, line_spacing=1.15):
    """Add multiple bullet lines in a single textbox."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txBox


def build_slide():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # ── Background ───────────────────────────────────────────────────────────
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    # ── Slide title ──────────────────────────────────────────────────────────
    add_text(slide, "Conclusion",
             left=Inches(0.5), top=Inches(0.2),
             width=Inches(12.33), height=Inches(0.55),
             font_size=28, bold=True, color=WHITE,
             align=PP_ALIGN.LEFT)

    # ── Accent separator ─────────────────────────────────────────────────────
    add_rect(slide,
             left=Inches(0.5), top=Inches(0.78),
             width=Inches(12.33), height=Pt(1.5),
             fill_color=ACCENT_LINE)

    # ── Central statement box ────────────────────────────────────────────────
    stmt_top = Inches(0.92)
    stmt_bg  = add_rect(slide,
                        left=Inches(0.5), top=stmt_top,
                        width=Inches(12.33), height=Inches(0.9),
                        fill_color=RGBColor(0x10, 0x4A, 0x6E))

    add_text(slide,
             "T3.4 delivers a working, validated, RAI-compliant blockchain "
             "governance system for heterogeneous MultiGIS data",
             left=Inches(0.65), top=stmt_top + Inches(0.08),
             width=Inches(12.0), height=Inches(0.75),
             font_size=14, bold=True, italic=True,
             color=WHITE, align=PP_ALIGN.CENTER)

    # ── Three takeaway blocks ─────────────────────────────────────────────────
    blocks = [
        {
            "accent": BLUE,
            "title":  "What we built",
            "lines": [
                "- 4-layer pilot-agnostic decentralised governance architecture",
                "- 3 Solidity smart contracts: registration, access control, provenance",
                "- Hybrid PostGIS / blockchain storage with SHA-256 hash bridge",
                "- Python LedgerInterface connecting pipeline to ledger",
            ],
        },
        {
            "accent": GREEN,
            "title":  "What we proved",
            "lines": [
                "- 8/8 components validated across 2 pilots",
                "- < 40 ms governance cycle — negligible overhead",
                "- 20/20 scalability — zero performance degradation",
                "- 7/7 fault tolerance — fully enforced at contract level",
            ],
        },
        {
            "accent": TEAL,
            "title":  "What it means for the project",
            "lines": [
                "- Trust layer for the entire WP3 data pipeline",
                "- Direct implementation of D2.2 architecture & D2.3 RAI requirements",
                "- Data foundation ready for WP4 AI models",
                "- Production migration to Hyperledger Besu — next phase",
            ],
        },
    ]

    BLOCK_TOP  = Inches(1.95)
    BLOCK_H    = Inches(4.6)
    BLOCK_W    = Inches(3.9)
    GAP        = Inches(0.27)
    LEFT_START = Inches(0.5)
    ACCENT_H   = Inches(0.06)
    TITLE_H    = Inches(0.5)
    LINES_TOP  = BLOCK_TOP + ACCENT_H + TITLE_H + Inches(0.15)
    LINES_H    = BLOCK_H - ACCENT_H - TITLE_H - Inches(0.25)

    for i, block in enumerate(blocks):
        left = LEFT_START + i * (BLOCK_W + GAP)

        # Card background
        add_rect(slide, left, BLOCK_TOP, BLOCK_W, BLOCK_H,
                 fill_color=CARD_BG)

        # Top accent bar
        add_rect(slide, left, BLOCK_TOP, BLOCK_W, ACCENT_H,
                 fill_color=block["accent"])

        # Block title
        add_text(slide, block["title"],
                 left=left + Inches(0.15),
                 top=BLOCK_TOP + ACCENT_H + Inches(0.1),
                 width=BLOCK_W - Inches(0.3),
                 height=TITLE_H,
                 font_size=14, bold=True,
                 color=block["accent"],
                 align=PP_ALIGN.LEFT)

        # Thin divider under title
        add_rect(slide,
                 left=left + Inches(0.15),
                 top=BLOCK_TOP + ACCENT_H + TITLE_H + Inches(0.08),
                 width=BLOCK_W - Inches(0.3),
                 height=Pt(0.8),
                 fill_color=RGBColor(0x17, 0x3A, 0x52))

        # Bullet lines
        add_text_lines(slide, block["lines"],
                       left=left + Inches(0.15),
                       top=LINES_TOP,
                       width=BLOCK_W - Inches(0.3),
                       height=LINES_H,
                       font_size=11,
                       color=RGBColor(0xCC, 0xD6, 0xDE))

    # ── Thank you line ───────────────────────────────────────────────────────
    add_text(slide,
             "Thank you  —  questions welcome",
             left=Inches(0.5), top=Inches(6.78),
             width=Inches(12.33), height=Inches(0.45),
             font_size=13, bold=True, italic=True,
             color=MID_GRAY, align=PP_ALIGN.CENTER)

    # ── Save ─────────────────────────────────────────────────────────────────
    out = "slide15_conclusion.pptx"
    prs.save(out)
    print(f"Saved -> {out}")


if __name__ == "__main__":
    build_slide()